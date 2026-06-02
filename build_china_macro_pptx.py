"""
把《中国企业生存与民众生活境况体检报告》转换为 PPTX
基于 generate_china_macro_report.py 的数据结构

输出：reports/China_Enterprise_Living_Diagnosis_2026-06-01.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os

from generate_china_macro_report import (
    REPORT_DATE, REPORT_TITLE, SUBTITLE,
    ENTERPRISE_DATA, PEOPLE_DATA, ZHENAI_IMPACT,
    STRATEGIC_PATHS, ZHAO_REVIEW,
)

# 颜色
RED = RGBColor(0xC0, 0x39, 0x2B)
RED_DARK = RGBColor(0x8E, 0x28, 0x20)
GREEN = RGBColor(0x00, 0xC8, 0x53)
GREEN_DARK = RGBColor(0x1E, 0x84, 0x49)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
GREY_DARK = RGBColor(0x2C, 0x3E, 0x50)
GREY = RGBColor(0x70, 0x70, 0x70)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

COLOR_BY_TAG = {"red": RED, "yellow": YELLOW, "green": GREEN}

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_font(run, size=14, bold=False, color=GREY_DARK, name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("eastAsia", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        el = rPr.makeelement(qn(f"a:{tag}"), {"typeface": name})
        rPr.append(el)


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=GREY_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)
    else:
        box.line.fill.background()
    return box


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_topbar_card(slide, x, y, w, h, color):
    """白色卡片 + 顶部彩条"""
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, x, y, w, Inches(0.05), color)


def add_footer(slide, page_no, total_pages):
    add_text(slide, Inches(0.4), Inches(7.18), Inches(8),
             Inches(0.25), "智慧助理 Linh · 报告日期 2026-06-01 · 数据来源：国家统计局/民政部/央行/市场监管总局/财新/长江商学院BCI/无破数据",
             size=9, color=GREY)
    add_text(slide, Inches(11.8), Inches(7.18), Inches(1.2),
             Inches(0.25), f"{page_no} / {total_pages}",
             size=9, color=GREY, align=PP_ALIGN.RIGHT)


def add_section_header(slide, num, title):
    add_rect(slide, Inches(0.4), Inches(0.35), Inches(0.7), Inches(0.7), RED)
    add_text(slide, Inches(0.4), Inches(0.35), Inches(0.7), Inches(0.7),
             num, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.25), Inches(0.4), Inches(11), Inches(0.65),
             title, size=24, bold=True, color=GREY_DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.4), Inches(1.13), Inches(12.5), Emu(20000), RED)


# =============================================================================
# Slide 1: 封面
# =============================================================================
def slide_cover(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # 背景渐变（用大矩形模拟）
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, RED_DARK)
    add_rect(s, 0, 0, SLIDE_W, Inches(5.5), RED)
    # 标题
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
             REPORT_TITLE, size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.6),
             SUBTITLE, size=20, color=WHITE)
    # 副信息条
    add_rect(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.2),
             RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, Inches(1.0), Inches(5.95), Inches(3.5), Inches(0.5),
             "📅 报告日期", size=12, color=GREY)
    add_text(s, Inches(1.0), Inches(6.35), Inches(3.5), Inches(0.5),
             REPORT_DATE.strftime("%Y-%m-%d"), size=18, bold=True, color=RED)
    add_text(s, Inches(4.5), Inches(5.95), Inches(3.5), Inches(0.5),
             "🧑‍💼 编制", size=12, color=GREY)
    add_text(s, Inches(4.5), Inches(6.35), Inches(3.5), Inches(0.5),
             "智慧助理 Linh", size=18, bold=True, color=RED)
    add_text(s, Inches(8.0), Inches(5.95), Inches(4.5), Inches(0.5),
             "👤 呈 / 抄送", size=12, color=GREY)
    add_text(s, Inches(8.0), Inches(6.35), Inches(4.5), Inches(0.5),
             "赵总  /  田小英", size=18, bold=True, color=RED)


# =============================================================================
# Slide 2: 执行摘要
# =============================================================================
def slide_exec_summary(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_section_header(s, "📌", "执行摘要 · 一页看完")

    points = [
        ("大盘 · 温和扩张但暗流涌动",
         "制造业PMI 50.3%、规上工业利润+18.2%看上去稳，但服务业PMI跌破荣枯线（49.4%）、长江商学院BCI仅46.9，反映「政府投资型行业向好、民营消费型行业承压」的鲜明分化。", YELLOW),
        ("居民 · 捂紧钱包+主动去杠杆",
         "Q1住户存款新增7.68万亿，但消费贷净减1,640亿。「赚得多花得少」，大额非必需消费（高客单婚恋）首当其冲。", RED),
        ("就业 · 青年群体仍是重灾区",
         "16-24岁失业率16.3%（高于上年同期15.8%），25-29岁7.4%，正是珍爱网核心客群，是高客单转化阻力的根本原因。", RED),
        ("婚恋赛道 · 2025反弹是脉冲，2026继续下行",
         "2025年676.3万对（+10.76%）是新《婚姻登记条例》刺激+疫情积压释放，2026 Q1立即回落-6.24%创新低。结构性下行不可逆。", RED),
        ("结构性机会 · 二婚客群池在扩",
         "2025协议离婚274.3万对（+12.2万），叠加诉讼离婚约90万，是高净值、高决策力、低投诉风险的优质客群。", GREEN),
        ("战略结论",
         "从「广撒网拉新」切换到「存量深耕+客群分层」，下一步首推路径A（公海盘活+中腰部补强），零额外预算、3个月见效。", GREEN_DARK),
    ]
    y = Inches(1.35)
    for i, (h, body, c) in enumerate(points):
        # 序号圆角
        add_rect(s, Inches(0.5), y, Inches(0.45), Inches(0.85), c)
        add_text(s, Inches(0.5), y, Inches(0.45), Inches(0.85),
                 str(i+1), size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 标题+正文
        add_text(s, Inches(1.05), y, Inches(11.8), Inches(0.32),
                 h, size=15, bold=True, color=c)
        add_text(s, Inches(1.05), y + Inches(0.32), Inches(11.8), Inches(0.55),
                 body, size=11, color=GREY_DARK)
        y += Inches(0.95)

    add_footer(s, page, total)


# =============================================================================
# Slide 3-4: 企业指标卡片（4列 x 1-2行）
# =============================================================================
def slide_indicator_grid(prs, page, total, sec_num, sec_title, items, cols=4):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_section_header(s, sec_num, sec_title)

    keys = list(items.keys())
    n = len(keys)
    rows = (n + cols - 1) // cols
    margin_x = Inches(0.4)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    available_w = SLIDE_W - margin_x * 2 - gap_x * (cols - 1)
    card_w = Emu(int(available_w / cols))
    available_h = Inches(5.85)
    card_h = Emu(int((available_h - gap_y * (rows - 1)) / rows))

    for i, k in enumerate(keys):
        d = items[k]
        col = i % cols
        row = i // cols
        x = margin_x + (card_w + gap_x) * col
        y = Inches(1.3) + (card_h + gap_y) * row
        c = COLOR_BY_TAG[d["color"]]

        add_topbar_card(s, x, y, card_w, card_h, c)
        # 内容
        pad = Inches(0.13)
        add_text(s, x + pad, y + Inches(0.12), card_w - pad * 2, Inches(0.32),
                 d["title"], size=10, bold=True, color=GREY)
        add_text(s, x + pad, y + Inches(0.42), card_w - pad * 2, Inches(0.55),
                 d["value"], size=22, bold=True, color=c)
        add_text(s, x + pad, y + Inches(0.97), card_w - pad * 2, Inches(0.28),
                 d["ref"], size=9, color=GREY)
        add_text(s, x + pad, y + Inches(1.22), card_w - pad * 2, Inches(0.28),
                 d["trend"], size=10, bold=True, color=GREY_DARK)
        # 解读
        verdict = d["verdict"].replace("<b>", "").replace("</b>", "")
        add_text(s, x + pad, y + Inches(1.50), card_w - pad * 2,
                 card_h - Inches(2.05),
                 verdict, size=9, color=GREY_DARK,
                 fill=RGBColor(0xFA, 0xFA, 0xFA))
        # 来源
        add_text(s, x + pad, y + card_h - Inches(0.45), card_w - pad * 2,
                 Inches(0.40),
                 f"📎 {d['source']}", size=8, color=GREY)

    add_footer(s, page, total)


# =============================================================================
# Slide: 传导映射表
# =============================================================================
def slide_impact(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_section_header(s, "三", "宏观 → 珍爱网传导映射")
    add_text(s, Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.3),
             "把上面的国家数据，翻译成对珍爱网经营的具体影响。", size=11, color=GREY)

    # 表头
    headers = [("权重", 0.9), ("宏观信号", 3.5), ("传导路径", 3.0), ("对珍爱网的具体影响", 5.4)]
    x = Inches(0.4)
    y = Inches(1.55)
    head_h = Inches(0.45)
    add_rect(s, x, y, Inches(12.8), head_h, GREY_DARK)
    cur = x
    for h, w in headers:
        add_text(s, cur, y, Inches(w), head_h, h,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur += Inches(w)

    # 行
    row_y = y + head_h
    row_h = Inches(0.65)
    for i, it in enumerate(ZHENAI_IMPACT):
        bg = WHITE if i % 2 == 0 else RGBColor(0xFA, 0xFB, 0xFC)
        add_rect(s, x, row_y, Inches(12.8), row_h, bg,
                 line=RGBColor(0xEE, 0xEE, 0xEE))
        cur = x
        cells = [
            (it["weight"], 0.9, RED, True, 12, PP_ALIGN.CENTER),
            (it["macro"], 3.5, GREY_DARK, True, 10, PP_ALIGN.LEFT),
            (it["channel"], 3.0, GREY_DARK, False, 10, PP_ALIGN.LEFT),
            (it["zhenai"], 5.4, RGBColor(0xB7, 0x1C, 0x1C), False, 10, PP_ALIGN.LEFT),
        ]
        for txt, w, col, bold, sz, al in cells:
            add_text(s, cur, row_y, Inches(w), row_h, txt,
                     size=sz, bold=bold, color=col, align=al,
                     anchor=MSO_ANCHOR.MIDDLE)
            cur += Inches(w)
        row_y += row_h

    add_footer(s, page, total)


# =============================================================================
# Slide: 战略路径单页
# =============================================================================
def slide_strategy_path(prs, page, total, idx):
    p = STRATEGIC_PATHS[idx]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    badge_color = [GREEN_DARK, YELLOW, RGBColor(0x21, 0x96, 0xF3)][idx]
    label = ["A · 存量深耕（推荐⭐⭐⭐⭐⭐）",
             "B · 低客单订阅突围（⭐⭐⭐⭐）",
             "C · 二婚高净值精耕（⭐⭐⭐⭐）"][idx]

    add_rect(s, Inches(0.4), Inches(0.35), Inches(0.7), Inches(0.7), badge_color)
    add_text(s, Inches(0.4), Inches(0.35), Inches(0.7), Inches(0.7),
             ["A", "B", "C"][idx], size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.25), Inches(0.4), Inches(11), Inches(0.65),
             f"四 · 战略路径 {label}", size=22, bold=True, color=GREY_DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.4), Inches(1.13), Inches(12.5), Emu(20000), badge_color)

    # 主张
    add_rect(s, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.7),
             RGBColor(0xE8, 0xF4, 0xFD))
    add_rect(s, Inches(0.4), Inches(1.35), Inches(0.06), Inches(0.7), badge_color)
    add_text(s, Inches(0.55), Inches(1.35), Inches(12.3), Inches(0.7),
             f"📌 核心主张：{p['thesis']}", size=14, bold=True, color=GREY_DARK,
             anchor=MSO_ANCHOR.MIDDLE)

    # 推演逻辑
    add_text(s, Inches(0.4), Inches(2.2), Inches(7), Inches(0.35),
             "🔍 推演逻辑", size=14, bold=True, color=GREY_DARK)
    box_logic = s.shapes.add_textbox(Inches(0.4), Inches(2.55), Inches(7.5), Inches(3.6))
    tf = box_logic.text_frame
    tf.word_wrap = True
    for j, line in enumerate(p["logic"]):
        para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = f"{j+1}. {line}"
        set_font(run, size=11, color=GREY_DARK)
        para.space_after = Pt(8)

    # 红线一致性
    add_text(s, Inches(8.1), Inches(2.2), Inches(5), Inches(0.35),
             "✅ 红线一致性自检", size=14, bold=True, color=GREEN_DARK)
    box_red = s.shapes.add_textbox(Inches(8.1), Inches(2.55), Inches(4.85), Inches(3.6))
    tf = box_red.text_frame
    tf.word_wrap = True
    for j, line in enumerate(p["fit_redlines"]):
        para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        set_font(run, size=11, color=GREY_DARK)
        para.space_after = Pt(6)

    # KPI 横条
    add_rect(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.65),
             RGBColor(0xFF, 0xF8, 0xE1), line=YELLOW)
    add_text(s, Inches(0.6), Inches(6.3), Inches(2.0), Inches(0.65),
             "💰 预估收益", size=11, bold=True, color=GREY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.5), Inches(6.3), Inches(6.4), Inches(0.65),
             p["expected"], size=12, bold=True, color=RED,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.1), Inches(6.3), Inches(1.5), Inches(0.65),
             "⚠️ 风险等级", size=11, bold=True, color=GREY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.5), Inches(6.3), Inches(2.4), Inches(0.65),
             p["risk"], size=12, bold=True, color=RED,
             anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


# =============================================================================
# Slide: 虚拟赵总质疑
# =============================================================================
def slide_zhao_review(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_section_header(s, "五", "虚拟赵总质疑 & 应对")
    add_text(s, Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.3),
             "重大方案输出前先自我质疑，逐条回应——保证合理性、可行性、红线一致性。", size=11, color=GREY)

    y = Inches(1.55)
    h_each = Inches(1.18)
    for i, (q, a) in enumerate(ZHAO_REVIEW["doubts"]):
        add_rect(s, Inches(0.4), y, Inches(12.5), h_each,
                 RGBColor(0xF9, 0xF9, 0xF9))
        add_rect(s, Inches(0.4), y, Inches(0.06), h_each, RED)
        add_text(s, Inches(0.6), y + Inches(0.08), Inches(12.2), Inches(0.45),
                 f"❓ 质疑 {i+1}：{q}", size=12, bold=True, color=RED)
        add_text(s, Inches(0.6), y + Inches(0.5), Inches(12.2), Inches(0.65),
                 f"✅ 应对：{a}", size=11, color=GREY_DARK)
        y += h_each + Inches(0.08)

    # 最终决断
    add_rect(s, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55),
             GREEN_DARK)
    add_text(s, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55),
             "📣 请赵总决断：建议 T+0 立即启动 路径A（存量深耕）；T+30 天根据 A 的成效在 B/C 中二选一加码。",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


# =============================================================================
# Slide: 数据来源与下次刷新
# =============================================================================
def slide_appendix(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_section_header(s, "📎", "数据来源清单 · 下次刷新")

    sources = [
        ("国家统计局",  "PMI / 工业利润 / 城镇调查失业率 / 收入消费 / 出生人口"),
        ("民政部",      "结婚登记 / 离婚登记 / 初婚年龄 / 季度民政统计数据"),
        ("中国人民银行","住户存款 / 消费贷 / 中长贷 / 一季度金融统计数据报告"),
        ("国家市场监督管理总局","新设经营主体 / 个体工商户 / 注销退出"),
        ("财新 / 标普全球","财新制造业PMI / 服务业PMI"),
        ("长江商学院","BCI 民营中小企业经营状况指数"),
        ("无破数据 / 全国企业破产重整案件信息网","破产案件总量、涉案企业、涉破资产"),
    ]
    y = Inches(1.45)
    for src, scope in sources:
        add_rect(s, Inches(0.4), y, Inches(0.06), Inches(0.6), RED)
        add_text(s, Inches(0.6), y, Inches(3.5), Inches(0.6),
                 src, size=13, bold=True, color=GREY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.2), y, Inches(8.7), Inches(0.6),
                 scope, size=11, color=GREY,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.7)

    # 下次刷新提示
    add_rect(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.65),
             RGBColor(0xFF, 0xF8, 0xE1), line=YELLOW)
    add_text(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.65),
             "🔁 下一次自动刷新预定：2026-07-01    |    一键刷新命令：python3 render_china_macro_report.py",
             size=12, bold=True, color=GREY_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


# =============================================================================
# 主程序
# =============================================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 计算总页数
    enterprise_pages = 1 if len(ENTERPRISE_DATA) <= 8 else 2
    people_pages = (len(PEOPLE_DATA) + 7) // 8  # 每页最多8张
    total_pages = (
        1  # cover
        + 1  # exec summary
        + enterprise_pages
        + people_pages
        + 1  # impact
        + 3  # 3 strategy paths
        + 1  # zhao review
        + 1  # appendix
    )

    page = 0
    slide_cover(prs); page += 1  # cover不计入页码
    slide_exec_summary(prs, 1, total_pages - 1)

    # 企业指标 8 项一页（4列 x 2行）
    slide_indicator_grid(prs, 2, total_pages - 1,
                         "一", "企业生存境况 · 8 项核心指标",
                         ENTERPRISE_DATA, cols=4)

    # 民生指标 11 项分两页：前 8 项 + 后 3 项
    keys = list(PEOPLE_DATA.keys())
    p1 = {k: PEOPLE_DATA[k] for k in keys[:8]}
    p2 = {k: PEOPLE_DATA[k] for k in keys[8:]}
    slide_indicator_grid(prs, 3, total_pages - 1,
                         "二", "民众生活境况 · 收入/消费/就业/储蓄（8项）",
                         p1, cols=4)
    slide_indicator_grid(prs, 4, total_pages - 1,
                         "二", "民众生活境况 · 婚恋人口（3项）",
                         p2, cols=3)

    slide_impact(prs, 5, total_pages - 1)
    slide_strategy_path(prs, 6, total_pages - 1, 0)
    slide_strategy_path(prs, 7, total_pages - 1, 1)
    slide_strategy_path(prs, 8, total_pages - 1, 2)
    slide_zhao_review(prs, 9, total_pages - 1)
    slide_appendix(prs, 10, total_pages - 1)

    out_path = f"reports/China_Enterprise_Living_Diagnosis_{REPORT_DATE.strftime('%Y-%m-%d')}.pptx"
    os.makedirs("reports", exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build()
    print(f"✅ PPT 已生成：{path}")
    print(f"   文件大小：{os.path.getsize(path)/1024:.1f} KB")
    print(f"   建议：浏览器或微信电脑端直接打开，PowerPoint/WPS/Keynote 均可")
